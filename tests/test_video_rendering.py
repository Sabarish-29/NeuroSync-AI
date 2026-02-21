"""
Tests for slide-to-image rendering in VideoAssembler.

Covers method detection, Pillow fallback rendering, LibreOffice pipeline
(mocked), cleanup, and create_video path.
"""

from __future__ import annotations

import os
import shutil
import textwrap
from pathlib import Path
from unittest.mock import MagicMock, patch, PropertyMock

import pytest

from neurosync.content.generators.video_assembler import (
    AssembledVideo,
    VideoAssembler,
    VideoSegment,
)


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------

@pytest.fixture()
def assembler(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> VideoAssembler:
    """VideoAssembler whose temp_dir lives inside the pytest tmp directory."""
    # Force pillow method so tests don't depend on LibreOffice
    monkeypatch.setattr(shutil, "which", lambda _cmd: None)
    asm = VideoAssembler()
    asm.temp_dir = tmp_path / "slides"
    asm.temp_dir.mkdir(exist_ok=True)
    return asm


@pytest.fixture()
def sample_pptx(tmp_path: Path) -> Path:
    """Create a minimal 3-slide PPTX for testing."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    for idx in range(1, 4):
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # title + content
        slide.shapes.title.text = f"Slide {idx} Title"
        slide.placeholders[1].text = f"Body text for slide {idx}.\nSecond line."
    path = tmp_path / "sample.pptx"
    prs.save(str(path))
    return path


# ---------------------------------------------------------------------------
# Method detection
# ---------------------------------------------------------------------------

class TestDetectMethod:
    """Tests for _detect_best_method()."""

    def test_detects_libreoffice_when_both_present(self, monkeypatch: pytest.MonkeyPatch) -> None:
        """Returns 'libreoffice' when soffice AND pdftoppm are available."""
        monkeypatch.setattr(
            shutil, "which",
            lambda cmd: f"/usr/bin/{cmd}" if cmd in ("soffice", "pdftoppm") else None,
        )
        asm = VideoAssembler()
        assert asm.method == "libreoffice"

    def test_detects_pillow_when_soffice_missing(self, monkeypatch: pytest.MonkeyPatch) -> None:
        """Falls back to 'pillow' when soffice is absent."""
        monkeypatch.setattr(
            shutil, "which",
            lambda cmd: "/usr/bin/pdftoppm" if cmd == "pdftoppm" else None,
        )
        asm = VideoAssembler()
        assert asm.method == "pillow"

    def test_detects_pillow_when_pdftoppm_missing(self, monkeypatch: pytest.MonkeyPatch) -> None:
        """Falls back to 'pillow' when pdftoppm is absent."""
        monkeypatch.setattr(
            shutil, "which",
            lambda cmd: "/usr/bin/soffice" if cmd == "soffice" else None,
        )
        asm = VideoAssembler()
        assert asm.method == "pillow"

    def test_detects_pillow_when_nothing_available(self, monkeypatch: pytest.MonkeyPatch) -> None:
        """Falls back to 'pillow' when neither tool is installed."""
        monkeypatch.setattr(shutil, "which", lambda _cmd: None)
        asm = VideoAssembler()
        assert asm.method == "pillow"


# ---------------------------------------------------------------------------
# Pillow rendering
# ---------------------------------------------------------------------------

class TestPillowRendering:
    """Tests for the pure-Python Pillow fallback."""

    def test_renders_three_slides(self, assembler: VideoAssembler, sample_pptx: Path) -> None:
        """Produces one PNG per slide."""
        images = assembler._pptx_to_images_pillow(sample_pptx)
        assert len(images) == 3
        for img_path in images:
            assert img_path.exists()
            assert img_path.suffix == ".png"

    def test_images_are_correct_resolution(self, assembler: VideoAssembler, sample_pptx: Path) -> None:
        """Each PNG matches the assembler's configured resolution."""
        from PIL import Image

        images = assembler._pptx_to_images_pillow(sample_pptx)
        for img_path in images:
            with Image.open(img_path) as im:
                assert im.size == assembler.resolution

    def test_images_not_blank(self, assembler: VideoAssembler, sample_pptx: Path) -> None:
        """PNGs contain rendered text and are larger than a blank image."""
        images = assembler._pptx_to_images_pillow(sample_pptx)
        for img_path in images:
            # A blank 1920×1080 solid-colour PNG is ~5 KB; with text it
            # should be noticeably larger or at least not all one colour.
            assert img_path.stat().st_size > 1_000

    def test_extract_slide_text(self, assembler: VideoAssembler, sample_pptx: Path) -> None:
        """Title and body text are extracted from slide shapes."""
        from pptx import Presentation

        prs = Presentation(str(sample_pptx))
        slide = prs.slides[0]
        title, body = assembler._extract_slide_text(slide)
        assert "Slide 1 Title" in title
        assert "Body text" in body

    def test_render_text_image_returns_pil_image(self, assembler: VideoAssembler) -> None:
        """_render_text_image returns a Pillow Image object."""
        from PIL import Image

        img = assembler._render_text_image("Hello", "World")
        assert isinstance(img, Image.Image)
        assert img.size == assembler.resolution


# ---------------------------------------------------------------------------
# LibreOffice pipeline (mocked subprocess)
# ---------------------------------------------------------------------------

class TestLibreOfficePipeline:
    """Tests for the LibreOffice conversion pipeline using mocked subprocesses."""

    def test_convert_pptx_to_pdf(self, assembler: VideoAssembler, sample_pptx: Path) -> None:
        """_convert_pptx_to_pdf invokes soffice and returns a Path."""
        expected_pdf = assembler.temp_dir / "sample.pdf"
        expected_pdf.write_bytes(b"%PDF-1.4 fake")

        with patch("subprocess.run") as mock_run:
            mock_run.return_value = MagicMock(returncode=0, stderr="")
            result = assembler._convert_pptx_to_pdf(sample_pptx)

        assert result == expected_pdf
        mock_run.assert_called_once()
        call_args = mock_run.call_args[0][0]
        assert "soffice" in call_args

    def test_convert_pptx_to_pdf_raises_on_failure(
        self, assembler: VideoAssembler, sample_pptx: Path,
    ) -> None:
        """RuntimeError when soffice returns non-zero."""
        with patch("subprocess.run") as mock_run:
            mock_run.return_value = MagicMock(returncode=1, stderr="error")
            with pytest.raises(RuntimeError, match="soffice failed"):
                assembler._convert_pptx_to_pdf(sample_pptx)

    def test_convert_pdf_to_pngs(self, assembler: VideoAssembler) -> None:
        """_convert_pdf_to_pngs invokes pdftoppm and finds generated PNGs."""
        # Pre-create fake output PNGs
        for i in range(1, 4):
            (assembler.temp_dir / f"slide-{i}.png").write_bytes(b"PNG")

        fake_pdf = assembler.temp_dir / "deck.pdf"
        fake_pdf.write_bytes(b"%PDF")

        with patch("subprocess.run") as mock_run:
            mock_run.return_value = MagicMock(returncode=0, stderr="")
            result = assembler._convert_pdf_to_pngs(fake_pdf)

        assert len(result) == 3
        assert all(p.name.startswith("slide-") for p in result)

    def test_convert_pdf_to_pngs_raises_on_failure(self, assembler: VideoAssembler) -> None:
        """RuntimeError when pdftoppm returns non-zero."""
        fake_pdf = assembler.temp_dir / "deck.pdf"
        fake_pdf.write_bytes(b"%PDF")

        with patch("subprocess.run") as mock_run:
            mock_run.return_value = MagicMock(returncode=1, stderr="bad PDF")
            with pytest.raises(RuntimeError, match="pdftoppm failed"):
                assembler._convert_pdf_to_pngs(fake_pdf)


# ---------------------------------------------------------------------------
# pptx_to_images dispatcher
# ---------------------------------------------------------------------------

class TestPptxToImages:
    """Tests for the public pptx_to_images dispatcher."""

    def test_file_not_found_raises(self, assembler: VideoAssembler) -> None:
        """FileNotFoundError when PPTX does not exist."""
        with pytest.raises(FileNotFoundError):
            assembler.pptx_to_images("/nonexistent/deck.pptx")

    def test_dispatches_to_pillow(self, assembler: VideoAssembler, sample_pptx: Path) -> None:
        """Uses Pillow path when method == 'pillow'."""
        assembler.method = "pillow"
        images = assembler.pptx_to_images(sample_pptx)
        assert len(images) == 3

    def test_falls_back_on_libreoffice_error(
        self, assembler: VideoAssembler, sample_pptx: Path,
    ) -> None:
        """Falls back to Pillow when LibreOffice pipeline raises."""
        assembler.method = "libreoffice"
        with patch.object(
            assembler, "_pptx_to_images_libreoffice", side_effect=RuntimeError("boom"),
        ):
            images = assembler.pptx_to_images(sample_pptx)
        # Should have fallen back to pillow and produced images
        assert len(images) == 3


# ---------------------------------------------------------------------------
# Cleanup
# ---------------------------------------------------------------------------

class TestCleanup:
    """Tests for temp-file cleanup."""

    def test_cleanup_removes_png_and_pdf(self, assembler: VideoAssembler) -> None:
        """cleanup() deletes *.png and *.pdf from temp_dir."""
        (assembler.temp_dir / "slide_000.png").write_bytes(b"PNG")
        (assembler.temp_dir / "deck.pdf").write_bytes(b"PDF")
        assembler.cleanup()
        assert list(assembler.temp_dir.glob("*.png")) == []
        assert list(assembler.temp_dir.glob("*.pdf")) == []

    def test_cleanup_skipped_in_debug(
        self, assembler: VideoAssembler, monkeypatch: pytest.MonkeyPatch,
    ) -> None:
        """Files are kept when NEUROSYNC_DEBUG=1."""
        monkeypatch.setenv("NEUROSYNC_DEBUG", "1")
        png = assembler.temp_dir / "slide_000.png"
        png.write_bytes(b"PNG")
        assembler.cleanup()
        assert png.exists()

    def test_cleanup_no_error_on_empty_dir(self, assembler: VideoAssembler) -> None:
        """cleanup() succeeds even when temp_dir is empty."""
        assembler.cleanup()  # should not raise

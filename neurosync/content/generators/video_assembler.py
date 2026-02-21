"""
NeuroSync AI — Video Assembler.

Assembles narrated video from slides, audio clips, and diagrams
using MoviePy.  Supports two slide-rendering backends:

* **libreoffice** — PPTX → PDF → PNG via ``soffice`` + ``pdftoppm``
  (highest fidelity, Linux / macOS).
* **pillow** — text-only rendering via Pillow (cross-platform fallback).

The best available method is auto-detected at construction time.
"""

from __future__ import annotations

import glob
import os
import shutil
import subprocess
import textwrap
import time
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Optional

from loguru import logger


# ---------------------------------------------------------------------------
# Data classes
# ---------------------------------------------------------------------------

@dataclass
class VideoSegment:
    """A single segment of the video (one slide with audio)."""
    slide_title: str
    audio_path: Optional[str] = None
    image_path: Optional[str] = None
    duration_seconds: float = 8.0
    text_overlay: str = ""


@dataclass
class AssembledVideo:
    """Result of video assembly."""
    output_path: str
    total_duration_seconds: float
    segment_count: int
    resolution: tuple[int, int] = (1920, 1080)
    fps: int = 24
    file_size_bytes: int = 0


# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

_BG_COLOR = (15, 16, 24)
_TITLE_FONT_SIZE = 72
_BODY_FONT_SIZE = 48
_TEXT_WRAP_WIDTH = 60
_LIBREOFFICE_TIMEOUT = 60   # seconds
_PDFTOPPM_TIMEOUT = 120     # seconds


# ---------------------------------------------------------------------------
# VideoAssembler
# ---------------------------------------------------------------------------

class VideoAssembler:
    """Assemble video from slides and audio using MoviePy.

    On construction the assembler probes the host for *LibreOffice* and
    *poppler-utils* (``pdftoppm``).  If both are present the high-fidelity
    ``libreoffice`` pipeline is used; otherwise the pure-Python ``pillow``
    fallback renders slide text onto blank images.
    """

    def __init__(
        self,
        fps: int = 24,
        resolution: tuple[int, int] = (1920, 1080),
        codec: str = "libx264",
        audio_codec: str = "aac",
        default_slide_duration: float = 8.0,
    ) -> None:
        self.fps = fps
        self.resolution = resolution
        self.codec = codec
        self.audio_codec = audio_codec
        self.default_slide_duration = default_slide_duration

        # Temp directory for intermediate artefacts
        self.temp_dir = Path("temp_slides")
        self.temp_dir.mkdir(exist_ok=True)

        # Auto-detect best rendering backend
        self.method = self._detect_best_method()
        logger.info(f"Video rendering method: {self.method}")

    # ------------------------------------------------------------------
    # Detection
    # ------------------------------------------------------------------

    def _detect_best_method(self) -> str:
        """Return ``'libreoffice'`` or ``'pillow'`` based on host tools."""
        has_soffice = shutil.which("soffice") is not None
        has_pdftoppm = shutil.which("pdftoppm") is not None

        if has_soffice and has_pdftoppm:
            return "libreoffice"

        logger.warning(
            "LibreOffice/poppler not found — falling back to Pillow "
            "text-only slide rendering"
        )
        return "pillow"

    # ------------------------------------------------------------------
    # PPTX → images (dispatcher)
    # ------------------------------------------------------------------

    def pptx_to_images(self, pptx_path: str | Path) -> list[Path]:
        """Convert a PPTX file to a list of PNG images.

        Automatically selects the best available pipeline.

        Args:
            pptx_path: Path to the ``.pptx`` file.

        Returns:
            Sorted list of ``Path`` objects pointing to generated PNGs.
        """
        pptx_path = Path(pptx_path)
        if not pptx_path.exists():
            raise FileNotFoundError(f"PPTX not found: {pptx_path}")

        if self.method == "libreoffice":
            try:
                return self._pptx_to_images_libreoffice(pptx_path)
            except Exception as exc:
                logger.error(f"LibreOffice pipeline failed: {exc}")
                logger.warning("Falling back to Pillow rendering")
                return self._pptx_to_images_pillow(pptx_path)

        return self._pptx_to_images_pillow(pptx_path)

    # ------------------------------------------------------------------
    # LibreOffice pipeline  (PPTX → PDF → PNG)
    # ------------------------------------------------------------------

    def _pptx_to_images_libreoffice(self, pptx_path: Path) -> list[Path]:
        """High-fidelity rendering via LibreOffice + pdftoppm."""
        pdf_path = self._convert_pptx_to_pdf(pptx_path)
        png_paths = self._convert_pdf_to_pngs(pdf_path)
        logger.info(f"Generated {len(png_paths)} slide images (libreoffice)")
        return png_paths

    def _convert_pptx_to_pdf(self, pptx_path: Path) -> Path:
        """Convert PPTX to PDF using ``soffice --headless``."""
        logger.info("Converting PPTX to PDF...")
        result = subprocess.run(
            [
                "soffice",
                "--headless",
                "--convert-to", "pdf",
                "--outdir", str(self.temp_dir),
                str(pptx_path),
            ],
            capture_output=True,
            text=True,
            timeout=_LIBREOFFICE_TIMEOUT,
        )
        if result.returncode != 0:
            raise RuntimeError(
                f"soffice failed (rc={result.returncode}): {result.stderr}"
            )

        pdf_name = pptx_path.stem + ".pdf"
        pdf_path = self.temp_dir / pdf_name
        if not pdf_path.exists():
            raise FileNotFoundError(
                f"Expected PDF not found after conversion: {pdf_path}"
            )
        return pdf_path

    def _convert_pdf_to_pngs(self, pdf_path: Path) -> list[Path]:
        """Convert PDF pages to PNG images using ``pdftoppm``."""
        logger.info("Converting PDF to PNG images...")
        output_prefix = self.temp_dir / "slide"

        result = subprocess.run(
            [
                "pdftoppm",
                "-png",
                "-r", "150",           # 150 DPI
                str(pdf_path),
                str(output_prefix),
            ],
            capture_output=True,
            text=True,
            timeout=_PDFTOPPM_TIMEOUT,
        )
        if result.returncode != 0:
            raise RuntimeError(
                f"pdftoppm failed (rc={result.returncode}): {result.stderr}"
            )

        # pdftoppm names files  slide-1.png, slide-2.png, …
        png_paths = sorted(self.temp_dir.glob("slide-*.png"))
        if not png_paths:
            raise FileNotFoundError(
                "pdftoppm produced no PNG files — check PDF content"
            )
        return png_paths

    # ------------------------------------------------------------------
    # Pillow fallback  (text-only rendering)
    # ------------------------------------------------------------------

    def _pptx_to_images_pillow(self, pptx_path: Path) -> list[Path]:
        """Render slide text onto 1920×1080 images with Pillow."""
        from PIL import Image, ImageDraw, ImageFont
        from pptx import Presentation

        logger.warning("Using basic text-only rendering")
        prs = Presentation(str(pptx_path))
        image_paths: list[Path] = []
        total = len(prs.slides)

        for i, slide in enumerate(prs.slides):
            t0 = time.monotonic()
            title, body = self._extract_slide_text(slide)
            img = self._render_text_image(title, body)

            out = self.temp_dir / f"slide_{i:03d}.png"
            img.save(str(out))
            image_paths.append(out)

            elapsed = time.monotonic() - t0
            logger.info(f"Processed slide {i + 1}/{total} ({elapsed:.1f}s)")

        logger.info(f"Generated {len(image_paths)} slide images (pillow)")
        return image_paths

    @staticmethod
    def _extract_slide_text(slide: Any) -> tuple[str, str]:
        """Pull title and body text from a python-pptx slide object."""
        title = ""
        body_parts: list[str] = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            # First shape with text becomes title if title is empty
            if not title:
                title = text
            else:
                body_parts.append(text)

        return title, "\n".join(body_parts)

    def _render_text_image(self, title: str, body: str) -> Any:
        """Create a 1920×1080 image with *title* and *body* text."""
        from PIL import Image, ImageDraw, ImageFont

        w, h = self.resolution
        img = Image.new("RGB", (w, h), color=_BG_COLOR)
        draw = ImageDraw.Draw(img)

        title_font = self._get_font(_TITLE_FONT_SIZE)
        body_font = self._get_font(_BODY_FONT_SIZE)

        # Title — centred near the top
        wrapped_title = textwrap.fill(title, width=_TEXT_WRAP_WIDTH)
        draw.multiline_text(
            (w // 2, 120),
            wrapped_title,
            font=title_font,
            fill="white",
            anchor="ma",
            align="center",
        )

        # Body — left-aligned below title
        if body:
            wrapped_body = textwrap.fill(body, width=_TEXT_WRAP_WIDTH)
            draw.multiline_text(
                (120, 320),
                wrapped_body,
                font=body_font,
                fill="white",
                align="left",
            )

        return img

    @staticmethod
    def _get_font(size: int) -> Any:
        """Return a TrueType font or the default bitmap font."""
        from PIL import ImageFont

        # Try common system font paths
        for candidate in (
            "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
            "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
            "/System/Library/Fonts/Helvetica.ttc",
            "C:\\Windows\\Fonts\\arial.ttf",
        ):
            if Path(candidate).exists():
                return ImageFont.truetype(candidate, size)

        return ImageFont.load_default()

    # ------------------------------------------------------------------
    # Segment creation
    # ------------------------------------------------------------------

    def create_segments(
        self,
        slide_titles: list[str],
        audio_paths: list[str | None] | None = None,
        durations: list[float] | None = None,
    ) -> list[VideoSegment]:
        """Create video segments from slide titles and optional audio.

        Args:
            slide_titles: List of slide titles/text.
            audio_paths: Optional list of audio file paths.
            durations: Optional list of segment durations.

        Returns:
            List of ``VideoSegment`` objects.
        """
        segments: list[VideoSegment] = []

        for i, title in enumerate(slide_titles):
            audio = audio_paths[i] if audio_paths and i < len(audio_paths) else None
            dur = durations[i] if durations and i < len(durations) else self.default_slide_duration

            segments.append(VideoSegment(
                slide_title=title,
                audio_path=audio,
                duration_seconds=dur,
                text_overlay=title,
            ))

        return segments

    # ------------------------------------------------------------------
    # Video assembly
    # ------------------------------------------------------------------

    def create_video(
        self,
        pptx_path: str | Path,
        audio_paths: list[str | Path],
        output_path: str | Path,
    ) -> AssembledVideo:
        """Build an MP4 video from a PPTX deck and per-slide audio files.

        Args:
            pptx_path: Path to the ``.pptx`` presentation.
            audio_paths: One audio file per slide.
            output_path: Destination ``.mp4`` path.

        Returns:
            ``AssembledVideo`` with metadata.
        """
        from moviepy import AudioFileClip, ImageClip, concatenate_videoclips

        image_paths = self.pptx_to_images(pptx_path)

        # Guard against slide / audio count mismatch
        n_slides = len(image_paths)
        n_audio = len(audio_paths)
        if n_slides != n_audio:
            logger.warning(
                f"Slide/audio count mismatch ({n_slides} slides, "
                f"{n_audio} audio files) — using min"
            )
        count = min(n_slides, n_audio)

        clips: list[Any] = []
        total_duration = 0.0
        path = Path(output_path)
        path.parent.mkdir(parents=True, exist_ok=True)

        try:
            for i in range(count):
                img = self._ensure_resolution(image_paths[i])
                audio = AudioFileClip(str(audio_paths[i]))
                clip = (
                    ImageClip(str(img))
                    .with_duration(audio.duration)
                    .with_audio(audio)
                )
                clips.append(clip)
                total_duration += audio.duration

            if not clips:
                raise ValueError("No segments to assemble")

            final = concatenate_videoclips(clips, method="compose")
            final.write_videofile(
                str(path),
                fps=self.fps,
                codec=self.codec,
                audio_codec=self.audio_codec,
                logger=None,
            )
        finally:
            # Free moviepy resources
            for c in clips:
                try:
                    c.close()
                except Exception:
                    pass

        file_size = path.stat().st_size if path.exists() else 0
        logger.info(
            f"Assembled video: {path.name} ({total_duration:.1f}s, "
            f"{count} segments, {file_size / (1024 * 1024):.1f}MB)"
        )

        return AssembledVideo(
            output_path=str(path),
            total_duration_seconds=total_duration,
            segment_count=count,
            resolution=self.resolution,
            fps=self.fps,
            file_size_bytes=file_size,
        )

    def assemble(self, segments: list[VideoSegment], output_path: str | Path) -> AssembledVideo:
        """Assemble video from pre-built segments (text-card mode).

        Creates text-card slides with optional audio overlay using
        MoviePy's ``TextClip`` for each slide and concatenates.

        Args:
            segments: List of video segments.
            output_path: Where to save the MP4.

        Returns:
            ``AssembledVideo`` result with metadata.
        """
        from moviepy import (
            ColorClip,
            CompositeVideoClip,
            TextClip,
            concatenate_videoclips,
        )

        path = Path(output_path)
        path.parent.mkdir(parents=True, exist_ok=True)

        clips = []
        total_duration = 0.0

        for seg in segments:
            duration = seg.duration_seconds or self.default_slide_duration

            bg = ColorClip(
                size=self.resolution,
                color=(30, 30, 60),
            ).with_duration(duration)

            try:
                txt = TextClip(
                    text=seg.slide_title,
                    font_size=48,
                    color="white",
                    size=(self.resolution[0] - 200, None),
                    method="caption",
                ).with_duration(duration).with_position("center")

                clip = CompositeVideoClip([bg, txt])
            except Exception:
                clip = bg

            clips.append(clip)
            total_duration += duration

        if not clips:
            raise ValueError("No segments to assemble")

        final = concatenate_videoclips(clips, method="compose")
        final.write_videofile(
            str(path),
            fps=self.fps,
            codec=self.codec,
            audio_codec=self.audio_codec,
            logger=None,
        )

        file_size = path.stat().st_size if path.exists() else 0
        logger.info(
            f"Assembled video: {path.name} ({total_duration:.1f}s, "
            f"{len(segments)} segments, {file_size / (1024 * 1024):.1f}MB)"
        )

        return AssembledVideo(
            output_path=str(path),
            total_duration_seconds=total_duration,
            segment_count=len(segments),
            resolution=self.resolution,
            fps=self.fps,
            file_size_bytes=file_size,
        )

    # ------------------------------------------------------------------
    # Helpers
    # ------------------------------------------------------------------

    def _ensure_resolution(self, image_path: Path) -> Path:
        """Resize an image to ``self.resolution`` if it doesn't match."""
        from PIL import Image

        with Image.open(image_path) as img:
            if img.size == self.resolution:
                return image_path
            resized = img.resize(self.resolution, Image.LANCZOS)
            resized.save(image_path)
        return image_path

    def estimate_duration(self, segments: list[VideoSegment]) -> float:
        """Estimate total video duration in seconds."""
        return sum(s.duration_seconds for s in segments)

    def cleanup(self) -> None:
        """Remove temporary slide images and intermediate PDFs.

        Skipped when the environment variable ``NEUROSYNC_DEBUG=1`` is set.
        """
        if os.environ.get("NEUROSYNC_DEBUG") == "1":
            logger.info("NEUROSYNC_DEBUG=1 — keeping temp files in {}", self.temp_dir)
            return

        removed = 0
        for pattern in ("*.png", "*.pdf"):
            for f in self.temp_dir.glob(pattern):
                f.unlink()
                removed += 1

        logger.info(f"Cleaned up {removed} temp files from {self.temp_dir}")
